import os, json, pandas as pd
from IPython.display import FileLink, display

# --- Notes: JSONL -> flattened CSV export ---
# Reads a JSONL file (one JSON object per line), flattens nested "data" fields,
# expands "solutions" into both a joined string and fixed columns (solution_1..solution_N),
# then writes a UTF-8-SIG CSV for Excel-friendly opening.

# --- Convert JSON results to flattened CSV ---
def jsonl_to_flattened_csv(
    json_path="agent_results_60.json",
    csv_path="agent_results_flattened.csv",
    max_solutions=10
):
    records = []

    if not os.path.exists(json_path):
        print("❌ JSON file not found.")
        return None

    with open(json_path, "r", encoding="utf-8") as f:
        for line in f:
            if not line.strip():
                continue

            row = json.loads(line)

            base = {
                "call_id": row.get("id"),
                "phone": row.get("phone"),
                "date": row.get("date"),
                "time": row.get("time"),
            }

            data = row.get("data", {})
            solutions = data.get("solutions", [])

            flat = {
                **base,
                "transcription": data.get("transcription"),
                "summary": data.get("summary"),
                "sentiment": data.get("sentiment"),
                "language_type": data.get("language_type"),
                "urgent_action": data.get("urgent_action"),
                "priority_score": data.get("priority_score"),
                "emotional_effect": data.get("emotional_effect"),
                "solutions": " | ".join(solutions),
            }

            for i in range(max_solutions):
                flat[f"solution_{i+1}"] = solutions[i] if i < len(solutions) else None

            records.append(flat)

    df = pd.DataFrame(records)
    df.to_csv(csv_path, index=False, encoding="utf-8-sig")

    return csv_path


csv_file = jsonl_to_flattened_csv()

if csv_file:
    print("✅ CSV created successfully")

import pandas as pd
import os

# --- Notes: Merge analysis output with demographics ---
# Loads the flattened agent CSV and customer demographics, normalizes phone to string,
# trims whitespace, drops missing phones, then performs an INNER JOIN on "phone".
# Result is written to merged.csv and downloaded in Colab.

# --- Merge agent_results_flattened.csv with customer_demographics.csv by "phone"---

# File paths
AGENT_CSV = "agent_results_flattened.csv"
DEMO_CSV = "customer_demographics.csv"
MERGED_CSV = "merged.csv"

# --- Safety checks ---
if not os.path.exists(AGENT_CSV):
    raise FileNotFoundError(f"❌ Missing file: {AGENT_CSV}")

if not os.path.exists(DEMO_CSV):
    raise FileNotFoundError(f"❌ Missing file: {DEMO_CSV}")

# --- Load CSVs ---
df_agent = pd.read_csv(AGENT_CSV, dtype={"phone": str})
df_demo = pd.read_csv(DEMO_CSV, dtype={"phone": str})

# --- Normalize phone numbers ---
df_agent["phone"] = df_agent["phone"].str.strip()
df_demo["phone"] = df_demo["phone"].str.strip()

# Optional: remove rows with missing phone
df_agent = df_agent[df_agent["phone"].notna()]
df_demo = df_demo[df_demo["phone"].notna()]

# --- INNER MERGE on phone ---
df_merged = pd.merge(
    df_agent,
    df_demo,
    on="phone",
    how="inner",
    suffixes=("_agent", "_demo")
)

# --- Save merged file ---
df_merged.to_csv(MERGED_CSV, index=False, encoding="utf-8-sig")

print(f"✅ Merged CSV created: {MERGED_CSV}")
print(f"📊 Rows in merged file: {len(df_merged)}")

# --- Download (Colab only) ---
from google.colab import files
files.download(MERGED_CSV)

print ("Done")

import pandas as pd
import numpy as np

# --- Notes: Feature engineering (age groups + income groups) ---
# Creates categorical segments for analysis:
# - ageGroup: bins age into 18-34, 35-54, 55+
# - incomeGroup2: collapses incomeGroup into "under $100k" vs "$100k+" (others -> NaN)
# Writes the enriched dataset to merged_with_ageGroup.csv.

# --- Create ageGroups and incomeGroup2 ---

# Load file
df = pd.read_csv("merged.csv")

# Ensure age is numeric
df["age"] = pd.to_numeric(df["age"], errors="coerce")

# Recode age into groups
df["ageGroup"] = pd.cut(
    df["age"],
    bins=[17, 34, 54, 100],
    labels=["18-34", "35-54", "55+"]
)

# Recode income into groups
def recode_income(val):
    if val in ["$50,000 to $74,999", "$75,000 to $99,999"]:
        return "under $100k"
    elif val in ["$100,000 to $124,999", "$125,000 to $149,999", "$150,000 or more"]:
        return "$100k+"
    else:
        return np.nan

df["incomeGroup2"] = df["incomeGroup"].apply(recode_income)

# Save updated file
df.to_csv("merged_with_ageGroup.csv", index=False, encoding="utf-8-sig")

print("✅ demo groups were created successfully")

import gradio as gr
import pandas as pd
import matplotlib.pyplot as plt
from statsmodels.stats.proportion import proportions_ztest

# ===========================================
# PART 2 – ANALYSIS DASHBOARD
# ===========================================
# Notes:
# - Loads merged_with_ageGroup.csv and derives time features (month/day/hour).
# - Builds plots:
#   1) Overall sentiment doughnut
#   2) Demographic sentiment bars with optional significance letters (z-tests)
#   3) Time trend line charts (month/day/hour)
# - Uses a dark theme to match the main app styling.

# ------------------------------
# Load CSV
# ------------------------------
df = pd.read_csv("merged_with_ageGroup.csv")
df["sentiment"] = df["sentiment"].str.strip().str.capitalize()

# ------------------------------
# Date & Time processing
# ------------------------------
# Notes:
# - Combines "date" + "time" into a single datetime.
# - errors="coerce" will convert invalid formats into NaT (missing datetime).
# - The date format here expects DD/MM/YYYY; ensure it matches what you saved earlier.

df["datetime"] = pd.to_datetime(
    df["date"] + " " + df["time"],
    format="%d/%m/%Y %H:%M:%S",
    errors="coerce"
)
df["month"] = df["datetime"].dt.month
# df["day_of_week"] = df["datetime"].dt.day_name()
df["day_of_week"] = df["datetime"].dt.strftime("%a") # abbreviation day name
df["hour"] = df["datetime"].dt.hour
df["minute"] = df["datetime"].dt.minute
df["month_name"] = df["datetime"].dt.month_name()
df["month_abbr"] = df["datetime"].dt.strftime("%b")

# ------------------------------
# Config
# ------------------------------
# Notes:
# - SENTIMENT_ORDER ensures consistent ordering across charts.
# - CATEGORY_ORDERS imposes stable category order for demographic axes.
# - TIME_LABELS provides axis label mapping.

SENTIMENT_COLORS = {
    "Positive": "#33cc33",
    "Neutral": "#ff9900",
    "Negative": "#ff3300"
}
SENTIMENT_ORDER = ["Negative", "Neutral", "Positive"]

BG_COLOR = "#2d2d2d"
AX_COLOR = "#3a3a3a"
TEXT_COLOR = "white"

GROUPING_VARS = ["gender", "ageGroup", "incomeGroup2", "education"]
CATEGORY_ORDERS = {
    "gender": ["Male", "Female"],
    "ageGroup": ["18-34", "35-54", "55+"],
    "incomeGroup2": ["under $100k", "$100k+"],
    "education": ["High School", "Bachelor", "Master", "PhD"]
}

VARIABLE_LABELS = {
    "gender": "Gender",
    "ageGroup": "Age Group",
    "incomeGroup2": "Income",
    "education": "Education"
}

MONTH_ABBR_ORDER = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                    "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
DOW_ORDER = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
HOUR_ORDER = list(range(26))

TIME_LABELS = {
    "month": "Month",
    "month_name": "Month",
    "month_abbr": "Month",
    "day_of_week": "Day of Week",
    "hour": "Hour"
}

# ------------------------------
# Helper: compute sentiment %
# ------------------------------
# Notes:
# - Produces a tidy table: [group_col, sentiment, percent].
# - Uses SENTIMENT_ORDER to keep columns stable and fill missing sentiments with 0.

def sentiment_percent(group_col):
    g = (
        df.groupby([group_col, "sentiment"])
        .size()
        .unstack()
        .reindex(columns=SENTIMENT_ORDER, fill_value=0)
    )
    g = g.div(g.sum(axis=1), axis=0) * 100
    g = g.reset_index().melt(
        id_vars=group_col,
        var_name="sentiment",
        value_name="percent"
    )
    return g

# ------------------------------
# Helper: significance letters
# ------------------------------
# Notes:
# - For each category within group_col, compares sentiment proportions pairwise.
# - Only assigns a letter to the "winning" sentiment when it is significantly larger than another sentiment.
# - Uses one-sample (within-group) z-tests with alternative="larger".
# - This is a heuristic annotation layer, not a full multiple-comparison correction framework.

def sentiment_significance_letters(group_col, alpha=0.05):
    results = {}
    sentiment_to_letter = {"Negative": "A", "Neutral": "B", "Positive": "C"}

    for group_value in df[group_col].dropna().unique():
        subset = df[df[group_col] == group_value]
        counts = {s: (subset["sentiment"] == s).sum() for s in SENTIMENT_ORDER}
        n = len(subset)
        letters = {s: [] for s in SENTIMENT_ORDER}

        for s1 in SENTIMENT_ORDER:
            for s2 in SENTIMENT_ORDER:
                if s1 == s2: continue
                if counts[s1] <= counts[s2]: continue
                stat, p = proportions_ztest([counts[s1], counts[s2]], [n, n], alternative="larger")
                if p < alpha:
                    letters[s1].append(sentiment_to_letter[s2])

        results[group_value] = {s: "".join(letters[s]) for s in SENTIMENT_ORDER}
    return results

# ------------------------------
# Bar chart: demographics
# ------------------------------
# Notes:
# - Plots grouped bars for Negative/Neutral/Positive percentages for each category.
# - Adds inline labels: "xx%" plus optional significance suffix letters from z-tests.
# - Uses a custom spacing/layout for readable clusters.

def plot_bar(group_col, figsize=(5.2,3.2)):
    data = sentiment_percent(group_col)
    sig_letters = sentiment_significance_letters(group_col)

    if group_col in CATEGORY_ORDERS:
        data[group_col] = pd.Categorical(
            data[group_col],
            categories=CATEGORY_ORDERS[group_col],
            ordered=True
        )

    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor(BG_COLOR)
    ax.set_facecolor(AX_COLOR)

    categories = data[group_col].cat.categories if group_col in CATEGORY_ORDERS else data[group_col].unique()
    n_groups = len(categories)
    n_bars = len(SENTIMENT_ORDER)
    bar_width = 0.25
    group_spacing = 0.1
    total_group_width = n_bars * bar_width

    bars = []
    for i, sentiment in enumerate(SENTIMENT_ORDER):
        y = (
            data[data["sentiment"] == sentiment]
            .set_index(group_col)
            .reindex(categories)["percent"]
            .fillna(0)
        )
        x_positions = [
            p * (total_group_width + group_spacing) + i * bar_width
            for p in range(n_groups)
        ]
        bar_container = ax.bar(
            x_positions,
            y,
            width=bar_width,
            color=SENTIMENT_COLORS[sentiment],
            label=sentiment
        )

        # Add percentage + significance letters (inline)
        for bar, p_idx in zip(bar_container, range(n_groups)):
            height = bar.get_height()
            group_value = categories[p_idx]
            suffix = sig_letters.get(group_value, {}).get(sentiment, "")

            label = f"{height:.0f}%"
            if suffix:
                label = f"{label} {suffix}"

            ax.text(
                bar.get_x() + bar.get_width() / 2,
                height + 0.8,
                label,
                ha="center",
                va="bottom",
                fontsize=8,
                color=TEXT_COLOR,
                fontweight="bold" if suffix else "normal"
            )

        bars.append(bar_container[0])

    # X-axis
    cluster_centers = [
        p * (total_group_width + group_spacing) + total_group_width / 2 - bar_width/2
        for p in range(n_groups)
    ]
    ax.set_xticks(cluster_centers)
    ax.set_xticklabels(categories, color=TEXT_COLOR)

    ax.set_ylim(0, 100)
    ax.set_ylabel("Percent (%)", color=TEXT_COLOR)
    ax.set_xlabel(VARIABLE_LABELS.get(group_col, group_col), color=TEXT_COLOR)
    ax.set_title(f"Sentiment by {VARIABLE_LABELS.get(group_col, group_col)}", color=TEXT_COLOR, pad=10)

    ax.tick_params(colors=TEXT_COLOR)
    for spine in ax.spines.values():
        spine.set_color("#555")

    ax.legend(handles=bars, facecolor=AX_COLOR, edgecolor="#555", labelcolor=TEXT_COLOR, fontsize=8)
    plt.tight_layout()
    plt.close(fig)  # prevents Matplotlib from auto-displaying
    return fig

# ------------------------------
# Line chart: time trends
# ------------------------------
# Notes:
# - Aggregates sentiment by a time dimension and normalizes to percentages by time bucket.
# - plot_sentiment_line draws one line per sentiment (Negative/Neutral/Positive).

def sentiment_over_time(time_col, order=None):
    g = (
        df.groupby([time_col, "sentiment"])
        .size()
        .unstack(fill_value=0)
        .reindex(columns=SENTIMENT_ORDER, fill_value=0)
    )
    if order is not None:
        g = g.reindex(order, fill_value=0)
        g.index.name = time_col
    g = g.div(g.sum(axis=1), axis=0).fillna(0) * 100
    g = g.reset_index()
    return g

def plot_sentiment_line(time_col, title, order=None, figsize=(5.5,3.5)):
    data = sentiment_over_time(time_col, order)
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor(BG_COLOR)
    ax.set_facecolor(AX_COLOR)

    for sentiment in SENTIMENT_ORDER:
        ax.plot(
            data[time_col],
            data[sentiment],
            marker="o",
            markersize=3,
            linewidth=1.5,
            label=sentiment,
            color=SENTIMENT_COLORS[sentiment]
        )

    ax.set_ylim(0,100)
    ax.set_ylabel("Percent (%)", color=TEXT_COLOR)
    ax.set_xlabel(TIME_LABELS.get(time_col,time_col), color=TEXT_COLOR)
    ax.set_title(title, color=TEXT_COLOR, pad=10)

    if time_col in ["month_name","month_abbr"]:
        ax.tick_params(axis="x", labelrotation=45, labelsize=8)
    elif time_col == "day_of_week":
        ax.set_xticks(range(len(DOW_ORDER)))
        ax.set_xticklabels(DOW_ORDER, color=TEXT_COLOR)
    elif time_col == "hour":
        ax.set_xticks(range(0,24,2))

    ax.tick_params(colors=TEXT_COLOR)
    for spine in ax.spines.values():
        spine.set_color("#555")
    ax.legend(facecolor=AX_COLOR, edgecolor="#555", labelcolor=TEXT_COLOR, fontsize=8)
    plt.tight_layout()
    plt.close(fig)  # prevents Matplotlib from auto-displaying
    return fig

# ------------------------------
# Doughnut pie chart
# ------------------------------
# Notes:
# - Uses a pie chart with wedge width < 1 to create a doughnut effect.
# - Percent labels are shown via autopct.

def plot_sentiment_pie(figsize=(4,4)):
    counts = df["sentiment"].value_counts().reindex(SENTIMENT_ORDER, fill_value=0)
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor(BG_COLOR)

    wedges, texts, autotexts = ax.pie(
        counts,
        labels=SENTIMENT_ORDER,
        autopct="%1.0f%%",
        colors=[SENTIMENT_COLORS[s] for s in SENTIMENT_ORDER],
        startangle=90,
        wedgeprops=dict(width=0.5)
    )
    for t in texts + autotexts:
        t.set_color(TEXT_COLOR)
        t.set_fontsize(9)
    ax.set_title("Overall Sentiment", color=TEXT_COLOR)
    plt.tight_layout()
    plt.close(fig)  # prevents Matplotlib from auto-displaying
    return fig

# ------------------------------
# Gradio dashboard
# ------------------------------
# Notes:
# - Static dashboard: plots are generated once at startup using value=...
# - If you want dynamic filters later, these plots would move into callbacks.

with gr.Blocks(css="""
    .gradio-container {
        background-color: #2d2d2d !important;
        color: white !important;
        font-family: Inter, sans-serif;
    }
""") as demo:

    gr.Markdown("## 📊 Sentiment Analytics Dashboard")

    # Top row: pie + bars
    with gr.Row():
        with gr.Column():
            gr.Plot(value=plot_sentiment_pie())
            gr.HTML("<div style='height:20px;'></div>")
            gr.Plot(value=plot_bar("gender"))
        with gr.Column():
            gr.Plot(value=plot_bar("ageGroup"))
            gr.HTML("<div style='height:20px;'></div>")
            gr.Plot(value=plot_bar("incomeGroup2"))

    gr.HTML("<hr style='border-color:#555'>")

    # Middle row: line charts
    with gr.Row():
        gr.Plot(value=plot_sentiment_line("month_abbr","Sentiment by Month", order=MONTH_ABBR_ORDER))
        gr.Plot(value=plot_sentiment_line("day_of_week","Sentiment by Day of Week", order=DOW_ORDER))
        gr.Plot(value=plot_sentiment_line("hour","Sentiment by Hour", order=HOUR_ORDER))

    gr.HTML("<hr style='border-color:#555'>")

demo.launch(share=True)

import pandas as pd  
import os
from typing import Dict
from pydantic import BaseModel, Field
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_openai import ChatOpenAI
from google.colab import userdata

# --- Notes: LLM summary layer (OpenAI via LangChain) ---
# This section defines:
# 1) A schema for aggregated sentiment stats + an executive summary
# 2) A parser enforcing JSON structure
# 3) A function that builds numeric distributions from df
# 4) A function that asks the model to write a short executive summary in English
# Requires OPENAI_API_KEY in Colab Secrets.


os.environ["OPENAI_API_KEY"] = userdata.get('OPENAI_API_KEY')

llm_instance = ChatOpenAI(
    model="gpt-4o",
    temperature=0
)

# --- 2. Pydantic Data Model ---
class SentimentSummarySchema(BaseModel):
    overall_sentiment: Dict[str, float] = Field(description="Overall sentiment percentages")
    sentiment_by_gender: Dict[str, Dict[str, float]] = Field(description="Sentiment percentages by gender")
    sentiment_by_ageGroup: Dict[str, Dict[str, float]] = Field(description="Sentiment percentages by age group")
    sentiment_by_incomeGroup2: Dict[str, Dict[str, float]] = Field(description="Sentiment percentages by income group")
    sentiment_by_education: Dict[str, Dict[str, float]] = Field(description="Sentiment percentages by education level")
    sentiment_by_month: Dict[str, Dict[str, float]] = Field(description="Monthly sentiment percentages")
    executive_summary: str = Field(description="Short AI-written executive summary (4–6 sentences)")

# --- 3. LangChain Parser ---
summary_parser = JsonOutputParser(pydantic_object=SentimentSummarySchema)

# --- 4. Prepare summary data ---
# Notes:
# - Converts raw rows into percentage distributions.
# - Output is a plain dict that is stringified before sending to the LLM.

def build_ai_summary_input(df: pd.DataFrame) -> dict:
    summary = {}
    summary["overall_sentiment"] = df["sentiment"].value_counts(normalize=True).mul(100).round(0).to_dict()

    for col in ["gender", "ageGroup", "incomeGroup2", "education"]:
        dist = df.groupby([col, "sentiment"]).size().unstack(fill_value=0)
        perc = dist.div(dist.sum(axis=1), axis=0).mul(100).round(0)
        summary[f"sentiment_by_{col}"] = perc.to_dict()

    monthly = df.groupby(["month_abbr", "sentiment"]).size().unstack(fill_value=0)
    monthly_perc = monthly.div(monthly.sum(axis=1), axis=0).mul(100).round(0)
    summary["sentiment_by_month"] = monthly_perc.to_dict()
    return summary

# --- 5. Analysis Logic ---
# Notes:
# - System prompt forces JSON-only English output following format_instructions.
# - If parsing fails, returns None and prints a diagnostic error.

def generate_ai_summary_langchain(df: pd.DataFrame):
    summary_data = build_ai_summary_input(df)

    prompt = ChatPromptTemplate.from_messages([
        ("system", "You are a senior data analyst. Write a structured JSON response in ENGLISH.\n{format_instructions}"),
        ("user", "Analyze the following sentiment data and produce an executive summary:\n{input_text}")
    ])

    
    chain = prompt | llm_instance | summary_parser

    try:
        return chain.invoke({
            "input_text": str(summary_data), # המרה למחרוזת כדי שה-LLM יוכל לקרוא
            "format_instructions": summary_parser.get_format_instructions()
        })
    except Exception as e:
        print(f"❌ Error generating AI summary: {e}")
        return None

# --- 6. Example usage ---
# Notes:
# - Placeholder main guard; actual execution occurs later when chaining to strategy generation.

if __name__ == "__main__":
    
    # ai_result = generate_ai_summary_langchain(df)
    pass

# =============================================
# PART 4 – STRATEGIC RECOMMENDATIONS USING AI
# =============================================
# Notes:
# - Takes the executive summary text and generates a structured set of CX/ops recommendations.
# - Output is enforced via Pydantic schema + JsonOutputParser.

# ----- Strategy & recommendations layer ----
# Staffing recommendations (when / why)
# Demographic-specific support needs
# Product / technical friction signals
# Operational priorities
# Short actionable next steps

from typing import Dict, List

# --- 1. Pydantic Data Model ---
class StrategicRecommendationsSchema(BaseModel):
    staffing_recommendations: List[str] = Field(
        description="Recommendations on staffing levels, shifts, and peak workloads"
    )
    demographic_support_insights: Dict[str, List[str]] = Field(
        description="Customer support needs by demographic group (e.g., gender, age)"
    )
    technical_friction_areas: List[str] = Field(
        description="Common technical or product-related difficulties"
    )
    operational_priorities: List[str] = Field(
        description="Top operational priorities to improve customer service"
    )
    executive_action_plan: str = Field(
        description="Concise 4–6 sentence executive action plan"
    )

# --- 2. LangChain Output Parser ---
strategy_parser = JsonOutputParser(
    pydantic_object=StrategicRecommendationsSchema
)

# --- 3. Strategic Analysis Logic ---
# Notes:
# - Uses the same llm_instance defined earlier (ChatOpenAI).
# - Returns JSON-only English recommendations.

def generate_strategic_recommendations(executive_summary: str):
    """
    Generates strategic customer service recommendations
    based on the executive summary from Part 2.
    """

    prompt = ChatPromptTemplate.from_messages([
        (
            "system",
            "You are a senior customer experience strategist. "
            "All outputs must be in ENGLISH and returned as valid JSON.\n"
            "{format_instructions}"
        ),
        (
            "user",
            "Based on the following executive summary, generate strategic "
            "recommendations to improve customer service operations.\n\n"
            "Focus on:\n"
            "- Staffing needs and peak workloads\n"
            "- Demographic groups requiring additional assistance\n"
            "- Technical or product friction during purchase or usage\n"
            "- Clear, actionable business recommendations\n\n"
            "Executive Summary:\n{input_text}"
        )
    ])

    chain = prompt | llm_instance | strategy_parser

    try:
        return chain.invoke({
            "input_text": executive_summary,
            "format_instructions": strategy_parser.get_format_instructions()
        })
    except Exception as e:
        print(f"❌ Error generating strategic recommendations: {e}")
        return None

# --- 4.  Usage (Chained from Part 2) ---
# Notes:
# - Runs summary generation first, then feeds summary_obj.executive_summary into strategy generation.
# - Produces printed outputs for quick inspection in notebook.

if __name__ == "__main__":

    # Assumes Part 2 already ran and produced summary_obj
    summary_result = generate_ai_summary_langchain(df)

    if summary_result:
        summary_obj = SentimentSummarySchema(**summary_result)

        strategy_result = generate_strategic_recommendations(
            summary_obj.executive_summary
        )

        if strategy_result:
            strategy = StrategicRecommendationsSchema(**strategy_result)

            print("\n📌 EXECUTIVE ACTION PLAN")
            print(strategy.executive_action_plan)

            print("\n👥 Staffing Recommendations")
            for s in strategy.staffing_recommendations:
                print("-", s)

            print("\n🎯 Demographic Support Insights")
            for k, v in strategy.demographic_support_insights.items():
                print(f"{k}:")
                for item in v:
                    print(" •", item)

# ----- Data Driven PPTx generator function ----
# Notes:
# - Saves dashboard charts to PNG images.
# - Creates a PowerPoint deck:
#   Slide 1: overall sentiment pie
#   Slide 2: demographic bar charts (3 side-by-side)
#   Slide 3: time trend line charts (stacked)
#   Slides 4+: executive summary split into bullets (4 per slide)
#   Next slides: strategic recommendations split into bullets (4 per slide)
# - Uses python-pptx for layout and formatting.

# Dashboard image
# For Executive summary and Strategic Recommendations:
# Split into bullets one sentence per bullet
# Font size 16
# Up to 4 bullets per slide (auto-split across multiple slides if more than 4)
# Number the slide title, e.g., "Executive Summary (1)", "Executive Summary (2)"

# !pip install python-pptx matplotlib

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches
import matplotlib.pyplot as plt
import math

# -------------------------------
# 1️⃣ Save charts as images
# -------------------------------
# Notes:
# - Calls your plotting functions and persists figures as PNGs for PowerPoint insertion.
# - bbox_inches='tight' helps reduce extra whitespace around charts.

def save_dashboard_images():
    imgs = {}

    # Pie chart
    pie_path = "chart_pie.png"
    plot_sentiment_pie().savefig(pie_path, bbox_inches='tight')
    imgs['pie'] = pie_path

    # Bar charts
    bar_paths = []
    for col in ["gender","ageGroup","incomeGroup2"]:
        bar_path = f"chart_bar_{col}.png"
        plot_bar(col).savefig(bar_path, bbox_inches='tight')
        bar_paths.append(bar_path)
    imgs['bars'] = bar_paths

    # Line charts
    line_charts = [
        ("month_abbr","Sentiment by Month", MONTH_ABBR_ORDER),
        ("day_of_week","Sentiment by Day of Week", DOW_ORDER),
        ("hour","Sentiment by Hour", HOUR_ORDER)
    ]
    line_paths = []
    for col,title,order in line_charts:
        line_path = f"chart_line_{col}.png"
        plot_sentiment_line(col,title,order).savefig(line_path, bbox_inches='tight')
        line_paths.append(line_path)
    imgs['lines'] = line_paths

    return imgs

# -------------------------------
# 2️⃣ Helper: split text into slides
# -------------------------------
# Notes:
# - Converts a list of bullet strings into chunks of size <= max_bullets for multi-slide layouts.

def split_text_to_slides(text_list, max_bullets=4):
    """Split list of sentences into list of chunks with max_bullets each."""
    return [text_list[i:i+max_bullets] for i in range(0, len(text_list), max_bullets)]

# -------------------------------
# 3️⃣ Generate PPTX with split slides
# -------------------------------
# Notes:
# - Creates image-only slides first (charts), then text slides for summary/recommendations.
# - Centers images horizontally using SLIDE_WIDTH and computed offsets.
# - Uses slide_layouts[5] for blank-ish chart slides and slide_layouts[1] for title+content text slides.

def generate_clean_pptx_numbered(summary_obj, strategy_obj, pptx_file="Agent_10x_Report_Numbered.pptx"):
    prs = Presentation()
    chart_imgs = save_dashboard_images()

    SLIDE_WIDTH = Inches(10)  # Default PowerPoint slide width is 10 inches

    # -------- Slide 1: Pie Chart --------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Overall Customer Sentiment"
    img_width, img_height = Inches(4.5), Inches(4.5)
    left = (SLIDE_WIDTH - img_width) / 2  # center horizontally
    slide.shapes.add_picture(chart_imgs['pie'], left, Inches(1.5), width=img_width, height=img_height)

    # -------- Slide 2: Demographic Bar Charts --------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Sentiment by Demographics"  # <-- Added title
    bar_paths = chart_imgs['bars']
    bar_width, bar_height = Inches(2.8), Inches(2.8)
    total_width = len(bar_paths) * bar_width
    start_left = (SLIDE_WIDTH - total_width) / 2  # center all bars
    for i, img in enumerate(bar_paths):
        slide.shapes.add_picture(img, start_left + i*bar_width, Inches(1.5), width=bar_width, height=bar_height)

    # -------- Slide 3: Over-Time Line Charts --------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Sentiment Over-Time"  # <-- Added title
    line_paths = chart_imgs['lines']
    line_width, line_height = Inches(6.5), Inches(2)
    for i, img in enumerate(line_paths):
        left = (SLIDE_WIDTH - line_width) / 2
        slide.shapes.add_picture(img, left, Inches(1.5 + i*2.3), width=line_width, height=line_height)

    # -------- Slides 4+: Executive Summary --------
    # Notes:
    # - Splits by ". " to approximate sentences; each sentence becomes a bullet.
    # - Titles are numbered for continuity across multiple slides.

    summary_lines = [s.strip() for s in summary_obj.executive_summary.split('. ') if s]
    summary_chunks = split_text_to_slides(summary_lines, max_bullets=4)

    for idx, chunk in enumerate(summary_chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Executive Summary ({idx+1})"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for line in chunk:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.level = 0

    # -------- Slides for Strategic Recommendations --------
    # Notes:
    # - Combines multiple recommendation sections into a single bullet stream.
    # - Splits into multi-slides (4 bullets per slide) with numbered titles.

    strategy_lines = []

    # Staffing
    strategy_lines.extend(strategy_obj.staffing_recommendations)
    # Demographics
    for group, items in strategy_obj.demographic_support_insights.items():
        strategy_lines.extend([f"{group}: {item}" for item in items])
    # Technical friction
    strategy_lines.extend(strategy_obj.technical_friction_areas)
    # Operational priorities
    strategy_lines.extend(strategy_obj.operational_priorities)

    strategy_chunks = split_text_to_slides(strategy_lines, max_bullets=4)
    for idx, chunk in enumerate(strategy_chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Strategic Recommendations ({idx+1})"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for line in chunk:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.level = 0

    prs.save(pptx_file)
    return pptx_file

# -------------------------------
# 4️⃣ Generate and download
# -------------------------------
# Notes:
# - Assumes summary_obj and strategy already exist in the notebook runtime.
# - Saves the PPTX and downloads it in Colab.

pptx_path = generate_clean_pptx_numbered(summary_obj, strategy)

from google.colab import files
files.download(pptx_path)
